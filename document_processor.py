import os
import hashlib
import re
import time
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed
from multiprocessing import cpu_count
from config import CHUNK_SIZE, CHUNK_OVERLAP
from config_loader import load_config

# Try to import tiktoken for token-based chunking
try:
    import tiktoken
    TIKTOKEN_AVAILABLE = True
except ImportError:
    TIKTOKEN_AVAILABLE = False
    print("[WARNING] tiktoken not available. Install with: pip install tiktoken")
    print("[INFO] Falling back to character-based chunking")

# Import Docling VLM processor with Granite integration
try:
    from pdf_processor_docling import DoclingPDFProcessor, pdf_page_to_multimodal_document
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("[WARNING] Docling with Granite VLM not available. PDF processing disabled.")

# Import Office document processors
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("[WARNING] python-pptx not available. PPTX processing disabled.")

try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("[WARNING] python-docx not available. DOCX processing disabled.")


# Global function for parallel processing (must be at module level for pickling)
def _process_single_file(args):
    """
    Process a single file. Used by ProcessPoolExecutor.
    
    Args:
        args: tuple of (file_path, directory_path, file_type)
        file_type: 'pdf', 'pptx', 'docx', or 'text'
    
    Returns:
        tuple of (document, relative_path, error)
    """
    file_path, directory_path, file_type = args
    rel_file = os.path.relpath(file_path, directory_path)
    
    try:
        if file_type == 'pdf':
            # Use Docling VLM pipeline with Granite API
            if not PDF_SUPPORT:
                raise ImportError("Docling with Granite VLM is required for PDF processing. Install with: pip install docling")
            
            config = load_config()
            pdf_config = config.get('document_processing', {}).get('pdf', {}).get('mode_2', {})
            
            # Extract hostname:port from api_url
            api_url_full = pdf_config.get('api_url', 'http://100.126.235.19:2222/v1/chat/completions')
            match = re.search(r'https?://([^/]+)', api_url_full)
            hostname_port = match.group(1) if match else "100.126.235.19:2222"
            
            # Initialize Docling VLM processor with Granite
            pdf_processor = DoclingPDFProcessor(
                api_url=hostname_port,
                model_name=pdf_config.get('model_name', 'ibm-granite-vl'),
                images_scale=pdf_config.get('images_scale', 2.0),
                timeout=pdf_config.get('timeout', 120),
                prompt=pdf_config.get('user_prompt', 'Convert this page to docling.')
            )
            
            # Process PDF
            pages = pdf_processor.process_pdf(file_path)
            
            # Combine text from all pages
            all_text = []
            for page in pages:
                page_header = f"\n\n=== Page {page.page_number} ===\n\n"
                all_text.append(page_header + page.text)
            
            content = "".join(all_text)
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path),
                'type': 'pdf',
                'total_pages': len(pages),
                'processor': 'granite_vlm'
            }
            
            doc = Document(content=content, metadata=metadata)
        elif file_type == 'pptx':
            # Process PPTX file
            if not PPTX_SUPPORT:
                raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
            
            presentation = Presentation(file_path)
            all_text = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_header = f"\n\n=== Slide {slide_idx} ===\n\n"
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    all_text.append(slide_header + "\n".join(slide_text))
            
            content = "".join(all_text) if all_text else ""
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path),
                'type': 'pptx',
                'total_slides': len(presentation.slides),
                'processor': 'python-pptx'
            }
            
            doc = Document(content=content, metadata=metadata)
        elif file_type == 'docx':
            # Process DOCX file
            if not DOCX_SUPPORT:
                raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
            
            docx_doc = DocxDocument(file_path)
            all_text = []
            
            # Extract text from paragraphs
            for para in docx_doc.paragraphs:
                if para.text.strip():
                    all_text.append(para.text.strip())
            
            # Extract text from tables
            for table in docx_doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    all_text.append("\n" + "\n".join(table_text) + "\n")
            
            content = "\n".join(all_text)
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path),
                'type': 'docx',
                'processor': 'python-docx'
            }
            
            doc = Document(content=content, metadata=metadata)
        else:
            # Regular text file
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path)
            }
            
            doc = Document(content=content, metadata=metadata)
        
        return (doc, rel_file, None)
    except Exception as e:
        return (None, rel_file, str(e))


class Document:
    """Represents a document with content and metadata."""
    
    def __init__(self, content: str, metadata: Dict[str, Any] = None):
        """
        Initialize a document.
        
        Args:
            content: The text content of the document
            metadata: Optional metadata dictionary
        """
        self.content = content
        self.metadata = metadata or {}
        
        # Generate SHA-256 content hash for duplicate detection
        self.content_hash = self._compute_content_hash(content)
        self.metadata['content_hash'] = self.content_hash
        
        # Add ingestion timestamp
        self.ingestion_timestamp = datetime.utcnow().isoformat() + 'Z'
        self.metadata['ingestion_timestamp'] = self.ingestion_timestamp
    
    @staticmethod
    def _compute_content_hash(content: str) -> str:
        """
        Compute SHA-256 hash of document content.
        
        Args:
            content: The document content
            
        Returns:
            Hexadecimal hash string
        """
        return hashlib.sha256(content.encode('utf-8')).hexdigest()
    
    def __repr__(self):
        return f"Document(content_length={len(self.content)}, metadata={self.metadata})"


class DocumentChunk:
    """Represents a chunk of a document."""
    
    def __init__(self, content: str, metadata: Dict[str, Any] = None, chunk_id: int = None):
        """
        Initialize a document chunk.
        
        Args:
            content: The text content of the chunk
            metadata: Optional metadata dictionary
            chunk_id: Unique identifier for the chunk
        """
        self.content = content
        self.metadata = metadata or {}
        self.chunk_id = chunk_id
    
    def __repr__(self):
        return f"DocumentChunk(id={self.chunk_id}, content_length={len(self.content)})"


class DocumentProcessor:
    """Processes documents and splits them into chunks."""
    
    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        """
        Initialize the document processor.
        Uses Granite VLM for all PDF processing.
        Supports token-based chunking with per-type overrides.
        
        Args:
            chunk_size: Size of each chunk (characters if legacy mode, tokens if token mode)
            chunk_overlap: Overlap between chunks (characters if legacy mode, tokens if token mode)
        """
        # Load chunking configuration
        config = load_config()
        chunking_config = config.get('document_processing', {}).get('chunking', {})
        
        # Determine chunking mode
        self.chunking_mode = chunking_config.get('mode', 'tokens')
        self.use_tokens = self.chunking_mode == 'tokens' and TIKTOKEN_AVAILABLE
        
        if self.use_tokens:
            # Token-based chunking
            tokenizer_model = chunking_config.get('tokenizer_model', 'cl100k_base')
            try:
                self.tokenizer = tiktoken.get_encoding(tokenizer_model)
                print(f"  [INFO] Using token-based chunking with {tokenizer_model}")
            except Exception as e:
                print(f"  [WARNING] Failed to load tokenizer {tokenizer_model}: {e}")
                print(f"  [INFO] Falling back to character-based chunking")
                self.use_tokens = False
                self.tokenizer = None
        else:
            self.tokenizer = None
            if self.chunking_mode == 'tokens' and not TIKTOKEN_AVAILABLE:
                print(f"  [WARNING] Token-based chunking requested but tiktoken not available")
                print(f"  [INFO] Falling back to character-based chunking")
        
        # Load per-type overrides
        self.per_type_overrides = chunking_config.get('per_type_overrides', {})
        
        # Set default chunk sizes
        if self.use_tokens:
            self.chunk_size = chunk_size or chunking_config.get('chunk_size_tokens', 512)
            self.chunk_overlap = chunk_overlap or chunking_config.get('chunk_overlap_tokens', 50)
            self.min_chunk_size = chunking_config.get('min_chunk_size_tokens', 50)
            self.max_chunk_size = chunking_config.get('max_chunk_size_tokens', 2048)
        else:
            self.chunk_size = chunk_size or chunking_config.get('chunk_size', CHUNK_SIZE)
            self.chunk_overlap = chunk_overlap or chunking_config.get('chunk_overlap', CHUNK_OVERLAP)
            self.min_chunk_size = chunking_config.get('min_chunk_size', 50)
            self.max_chunk_size = chunking_config.get('max_chunk_size', 2000)
        
        self.respect_sentence_boundary = chunking_config.get('respect_sentence_boundary', True)
        self.respect_paragraph_boundary = chunking_config.get('respect_paragraph_boundary', True)
        
        # Initialize PDF processor (Granite VLM only)
        self.pdf_processor = None
        self.processor_name = 'granite_vlm'
        
        if PDF_SUPPORT:
            # Load configuration for Granite VLM
            pdf_config = config.get('document_processing', {}).get('pdf', {}).get('mode_2', {})
            
            # Extract hostname:port from api_url
            api_url_full = pdf_config.get('api_url', 'http://100.126.235.19:2222/v1/chat/completions')
            match = re.search(r'https?://([^/]+)', api_url_full)
            hostname_port = match.group(1) if match else "100.126.235.19:2222"
            
            # Initialize Docling VLM processor with Granite
            self.pdf_processor = DoclingPDFProcessor(
                api_url=hostname_port,
                model_name=pdf_config.get('model_name', 'ibm-granite-vl'),
                images_scale=pdf_config.get('images_scale', 2.0),
                timeout=pdf_config.get('timeout', 120),
                prompt=pdf_config.get('user_prompt', 'Convert this page to docling.')
            )
        else:
            print("  [WARNING] Docling with Granite VLM not available")
            print("  [INFO] Install dependencies: pip install docling")
    
    def _detect_content_type(self, document: Document) -> str:
        """
        Detect content type from document metadata and content.
        
        Args:
            document: Document to analyze
            
        Returns:
            Content type string: 'code', 'table', 'documentation', 'pdf', 'office', 'structured', or 'default'
        """
        # Check file extension from metadata
        source = document.metadata.get('source', '')
        filename = document.metadata.get('filename', '')
        file_path = source or filename
        
        if file_path:
            file_lower = file_path.lower()
            
            # Check per-type overrides by extension
            for type_name, type_config in self.per_type_overrides.items():
                extensions = type_config.get('extensions', [])
                if any(file_lower.endswith(ext.lower()) for ext in extensions):
                    return type_name
        
        # Check content patterns for table detection
        content = document.content
        if self.per_type_overrides.get('table'):
            patterns = self.per_type_overrides['table'].get('patterns', [])
            for pattern in patterns:
                if pattern in content:
                    # Additional check: count pattern occurrences
                    if content.count(pattern) > 10:  # Likely a table
                        return 'table'
        
        # Default type
        return 'default'
    
    def _get_chunk_size_for_type(self, content_type: str) -> Tuple[int, int]:
        """
        Get chunk size and overlap for a specific content type.
        
        Args:
            content_type: Content type string
            
        Returns:
            Tuple of (chunk_size, chunk_overlap)
        """
        if content_type in self.per_type_overrides:
            type_config = self.per_type_overrides[content_type]
            if self.use_tokens:
                chunk_size = type_config.get('chunk_size_tokens', self.chunk_size)
                chunk_overlap = type_config.get('chunk_overlap_tokens', self.chunk_overlap)
            else:
                # Fallback to character-based if tokens not available
                chunk_size = type_config.get('chunk_size', self.chunk_size)
                chunk_overlap = type_config.get('chunk_overlap', self.chunk_overlap)
            return (chunk_size, chunk_overlap)
        
        return (self.chunk_size, self.chunk_overlap)
    
    def load_text_file(self, file_path: str) -> Document:
        """
        Load a text file as a document.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            Document object
        """
        file_lower = file_path.lower()
        
        # Handle PDF files separately
        if file_lower.endswith('.pdf'):
            return self._load_pdf_file(file_path)
        
        # Handle PPTX files
        elif file_lower.endswith('.pptx'):
            return self._load_pptx_file(file_path)
        
        # Handle DOCX files
        elif file_lower.endswith('.docx'):
            return self._load_docx_file(file_path)
        
        # Regular text file
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        metadata = {
            'source': file_path,
            'filename': os.path.basename(file_path)
        }
        
        return Document(content=content, metadata=metadata)
    
    def _load_pdf_file(self, file_path: str) -> Document:
        """
        Load a PDF file as a document using Docling or PyMuPDF.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Document object with extracted text from all pages
        """
        if not PDF_SUPPORT or not self.pdf_processor:
            raise ImportError(
                "PDF support not available. Install docling: pip install docling"
            )
        
        try:
            # Process PDF pages
            pages = self.pdf_processor.process_pdf(file_path)
            
            # Combine text from all pages
            all_text = []
            for page in pages:
                page_header = f"\n\n=== Page {page.page_number} ===\n\n"
                all_text.append(page_header + page.text)
            
            content = "".join(all_text)
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path),
                'type': 'pdf',
                'total_pages': len(pages),
                'processor': self.processor_name
            }
            
            return Document(content=content, metadata=metadata)
            
        except Exception as e:
            # If PDF processing fails, raise with descriptive error
            raise Exception(f"Failed to process PDF {file_path}: {str(e)}")
    
    def _load_pptx_file(self, file_path: str) -> Document:
        """
        Load a PPTX file as a document.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Document object with extracted text from all slides
        """
        if not PPTX_SUPPORT:
            raise ImportError(
                "PPTX support not available. Install python-pptx: pip install python-pptx"
            )
        
        try:
            presentation = Presentation(file_path)
            all_text = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_header = f"\n\n=== Slide {slide_idx} ===\n\n"
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    all_text.append(slide_header + "\n".join(slide_text))
            
            content = "".join(all_text) if all_text else ""
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path),
                'type': 'pptx',
                'total_slides': len(presentation.slides),
                'processor': 'python-pptx'
            }
            
            return Document(content=content, metadata=metadata)
            
        except Exception as e:
            raise Exception(f"Failed to process PPTX {file_path}: {str(e)}")
    
    def _load_docx_file(self, file_path: str) -> Document:
        """
        Load a DOCX file as a document.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Document object with extracted text from paragraphs and tables
        """
        if not DOCX_SUPPORT:
            raise ImportError(
                "DOCX support not available. Install python-docx: pip install python-docx"
            )
        
        try:
            docx_doc = DocxDocument(file_path)
            all_text = []
            
            # Extract text from paragraphs
            for para in docx_doc.paragraphs:
                if para.text.strip():
                    all_text.append(para.text.strip())
            
            # Extract text from tables
            for table in docx_doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    all_text.append("\n" + "\n".join(table_text) + "\n")
            
            content = "\n".join(all_text)
            
            metadata = {
                'source': file_path,
                'filename': os.path.basename(file_path),
                'type': 'docx',
                'processor': 'python-docx'
            }
            
            return Document(content=content, metadata=metadata)
            
        except Exception as e:
            raise Exception(f"Failed to process DOCX {file_path}: {str(e)}")
    
    def _load_file_safe(self, file_path: str, directory_path: str) -> tuple:
        """
        Safely load a file and return result with status.
        Helper method for parallel processing.
        """
        rel_file = os.path.relpath(file_path, directory_path)
        try:
            doc = self.load_text_file(file_path)
            return (doc, rel_file, None)  # (document, path, error)
        except Exception as e:
            return (None, rel_file, str(e))
    
    def load_directory(self, directory_path: str, extensions: List[str] = None, recursive: bool = True, parallel: bool = None, max_workers: int = None) -> List[Document]:
        """
        Load all text files from a directory.
        
        Args:
            directory_path: Path to the directory
            extensions: List of file extensions to include (e.g., ['.txt', '.md'])
            recursive: If True, search subdirectories recursively (default: True)
            parallel: If True, process files in parallel (None = use config setting)
            max_workers: Number of parallel workers (None = use config setting or CPU count)
            
        Returns:
            List of Document objects
        """
        if extensions is None:
            extensions = ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.pdf', '.pptx', '.docx']
        
        # Load config for parallel settings
        config = load_config()
        parallel_config = config.get('document_processing', {}).get('parallel', {})
        
        if parallel is None:
            parallel = parallel_config.get('enabled', True)
        
        if max_workers is None:
            max_workers = parallel_config.get('max_workers') or cpu_count()
        
        min_files_for_parallel = parallel_config.get('min_files_for_parallel', 2)
        parallel_mode = parallel_config.get('mode', 'process')  # 'process' or 'thread'
        
        documents = []
        
        # Collect all files first
        all_files = []
        if recursive:
            for root, _, files in os.walk(directory_path):
                for file in files:
                    if any(file.endswith(ext) for ext in extensions):
                        file_path = os.path.join(root, file)
                        all_files.append(file_path)
        else:
            try:
                files = os.listdir(directory_path)
                for file in files:
                    file_path = os.path.join(directory_path, file)
                    if os.path.isfile(file_path) and any(file.endswith(ext) for ext in extensions):
                        all_files.append(file_path)
            except Exception as e:
                print(f"Error accessing directory: {str(e)}")
                return documents
        
        if not all_files:
            return documents
        
        print(f"  Found {len(all_files)} files to process...")
        
        # Auto-detect I/O-heavy operations (PDFs with API calls benefit from threading)
        has_pdfs = any(f.lower().endswith('.pdf') for f in all_files)
        
        # For I/O-heavy operations (PDFs with API calls), prefer threading
        # For CPU-bound operations, prefer multiprocessing
        if parallel_mode == 'process' and has_pdfs:
            # Check if we should use threads for PDFs (I/O-heavy with API calls)
            io_heavy_config = parallel_config.get('use_threads_for_io', True)
            if io_heavy_config:
                print(f"  Detected I/O-heavy PDF processing - using ThreadPoolExecutor for better performance")
                parallel_mode = 'thread'
        
        if parallel and len(all_files) >= min_files_for_parallel:
            if parallel_mode == 'process':
                print(f"  Using {max_workers} parallel workers (multiprocessing - TRUE PARALLEL)")
                
                # Prepare arguments for parallel processing
                # Determine file type for each file
                def _get_file_type(file_path):
                    file_lower = file_path.lower()
                    if file_lower.endswith('.pdf'):
                        return 'pdf'
                    elif file_lower.endswith('.pptx'):
                        return 'pptx'
                    elif file_lower.endswith('.docx'):
                        return 'docx'
                    else:
                        return 'text'
                
                process_args = [
                    (file_path, directory_path, _get_file_type(file_path))
                    for file_path in all_files
                ]
                
                # Use ProcessPoolExecutor for true parallel processing (bypasses GIL)
                try:
                    with ProcessPoolExecutor(max_workers=max_workers) as executor:
                        # Submit all tasks
                        futures = [executor.submit(_process_single_file, args) for args in process_args]
                        
                        # Collect results as they complete
                        completed = 0
                        for future in as_completed(futures):
                            doc, rel_file, error = future.result()
                            completed += 1
                            
                            if doc:
                                documents.append(doc)
                                print(f"    [{completed}/{len(all_files)}] [OK] {rel_file}")
                            else:
                                print(f"    [{completed}/{len(all_files)}] [FAIL] {rel_file}: {error}")
                except Exception as e:
                    print(f"  [WARNING] Multiprocessing failed: {e}")
                    print(f"  [INFO] Falling back to sequential processing...")
                    parallel = False
            else:
                # Thread-based parallelism (better for I/O-heavy operations like PDF API calls)
                print(f"  Using {max_workers} parallel workers (threading - optimal for I/O-heavy operations)")
                
                # Get rate limiting settings
                worker_start_delay = parallel_config.get('worker_start_delay', 0.3)
                
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    future_to_file = {}
                    
                    # Submit tasks with delay to prevent overwhelming the API
                    for file_path in all_files:
                        future = executor.submit(self._load_file_safe, file_path, directory_path)
                        future_to_file[future] = file_path
                        # Small delay between submitting each task
                        time.sleep(worker_start_delay)
                    
                    completed = 0
                    for future in as_completed(future_to_file):
                        doc, rel_file, error = future.result()
                        completed += 1
                        
                        if doc:
                            documents.append(doc)
                            print(f"    [{completed}/{len(all_files)}] [OK] {rel_file}")
                        else:
                            print(f"    [{completed}/{len(all_files)}] [FAIL] {rel_file}: {error}")
        
        if not parallel or len(all_files) < min_files_for_parallel:
            # Sequential processing (for small file counts or when parallel is disabled)
            print(f"  Using sequential processing")
            for i, file_path in enumerate(all_files, 1):
                rel_file = os.path.relpath(file_path, directory_path)
                try:
                    doc = self.load_text_file(file_path)
                    documents.append(doc)
                    print(f"    [{i}/{len(all_files)}] [OK] {rel_file}")
                except Exception as e:
                    print(f"    [{i}/{len(all_files)}] [FAIL] {rel_file}: {str(e)}")
        
        return documents
    
    def split_text(self, text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
        """
        Split text into chunks with overlap.
        Uses token-based chunking if available, otherwise character-based.
        
        Args:
            text: The text to split
            chunk_size: Override chunk size (uses instance default if None)
            chunk_overlap: Override chunk overlap (uses instance default if None)
            
        Returns:
            List of text chunks
        """
        if chunk_size is None:
            chunk_size = self.chunk_size
        if chunk_overlap is None:
            chunk_overlap = self.chunk_overlap
        
        if self.use_tokens and self.tokenizer:
            return self._split_text_tokens(text, chunk_size, chunk_overlap)
        else:
            return self._split_text_characters(text, chunk_size, chunk_overlap)
    
    def _split_text_tokens(self, text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """
        Split text into chunks using token-based chunking.
        
        Args:
            text: The text to split
            chunk_size: Chunk size in tokens
            chunk_overlap: Overlap in tokens
            
        Returns:
            List of text chunks
        """
        if not text.strip():
            return []
        
        # Encode text to tokens
        tokens = self.tokenizer.encode(text)
        
        if len(tokens) <= chunk_size:
            return [text] if text.strip() else []
        
        chunks = []
        start_idx = 0
        
        while start_idx < len(tokens):
            end_idx = min(start_idx + chunk_size, len(tokens))
            chunk_tokens = tokens[start_idx:end_idx]
            
            # Decode tokens back to text
            chunk_text = self.tokenizer.decode(chunk_tokens)
            
            # If not the last chunk, try to break at sentence or word boundary
            if end_idx < len(tokens) and self.respect_sentence_boundary:
                # Try to find a sentence boundary within the last 25% of the chunk
                search_start = max(0, len(chunk_text) - chunk_size // 4)
                last_sentence = max(
                    chunk_text.rfind('. ', search_start),
                    chunk_text.rfind('? ', search_start),
                    chunk_text.rfind('! ', search_start),
                    chunk_text.rfind('.\n', search_start),
                    chunk_text.rfind('?\n', search_start),
                    chunk_text.rfind('!\n', search_start)
                )
                
                if last_sentence > len(chunk_text) // 2:
                    # Re-encode to find the token position
                    truncated_text = chunk_text[:last_sentence + 1]
                    truncated_tokens = self.tokenizer.encode(truncated_text)
                    if len(truncated_tokens) >= chunk_size // 2:
                        chunk_text = truncated_text
                        end_idx = start_idx + len(truncated_tokens)
                elif ' ' in chunk_text[search_start:]:
                    # Try to break at word boundary
                    last_space = chunk_text.rfind(' ', search_start)
                    if last_space > len(chunk_text) // 2:
                        truncated_text = chunk_text[:last_space]
                        truncated_tokens = self.tokenizer.encode(truncated_text)
                        if len(truncated_tokens) >= chunk_size // 2:
                            chunk_text = truncated_text
                            end_idx = start_idx + len(truncated_tokens)
            
            chunk_text = chunk_text.strip()
            if chunk_text and len(self.tokenizer.encode(chunk_text)) >= self.min_chunk_size:
                chunks.append(chunk_text)
            
            # Move start position with overlap
            if end_idx >= len(tokens):
                break
            
            # Calculate overlap in tokens
            overlap_tokens = min(chunk_overlap, len(chunk_tokens))
            start_idx = end_idx - overlap_tokens
            
            # Avoid infinite loop
            if start_idx >= end_idx:
                start_idx = end_idx
        
        return chunks
    
    def _split_text_characters(self, text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """
        Split text into chunks using character-based chunking (legacy method).
        
        Args:
            text: The text to split
            chunk_size: Chunk size in characters
            chunk_overlap: Overlap in characters
            
        Returns:
            List of text chunks
        """
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # If not the last chunk, try to break at a sentence or word boundary
            if end < len(text) and self.respect_sentence_boundary:
                # Try to find the last period, question mark, or exclamation point
                last_sentence = max(
                    chunk.rfind('. '),
                    chunk.rfind('? '),
                    chunk.rfind('! '),
                    chunk.rfind('.\n'),
                    chunk.rfind('?\n'),
                    chunk.rfind('!\n')
                )
                
                if last_sentence > chunk_size // 2:
                    chunk = chunk[:last_sentence + 1]
                    end = start + last_sentence + 1
                # Otherwise, try to find the last space
                elif ' ' in chunk:
                    last_space = chunk.rfind(' ')
                    if last_space > chunk_size // 2:
                        chunk = chunk[:last_space]
                        end = start + last_space
            
            chunk = chunk.strip()
            if chunk and len(chunk) >= self.min_chunk_size:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - chunk_overlap
            
            # Avoid infinite loop
            if start <= end - chunk_size:
                start = end
        
        return chunks
    
    def chunk_document(self, document: Document) -> List[DocumentChunk]:
        """
        Split a document into chunks with per-type overrides.
        
        Args:
            document: Document to chunk
            
        Returns:
            List of DocumentChunk objects
        """
        # Detect content type and get appropriate chunk sizes
        content_type = self._detect_content_type(document)
        chunk_size, chunk_overlap = self._get_chunk_size_for_type(content_type)
        
        # Add content type to metadata
        document.metadata['content_type'] = content_type
        
        # Split text using appropriate chunk sizes
        text_chunks = self.split_text(document.content, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        
        chunks = []
        for i, chunk_text in enumerate(text_chunks):
            metadata = document.metadata.copy()
            metadata['chunk_index'] = i
            metadata['total_chunks'] = len(text_chunks)
            metadata['chunk_size_used'] = chunk_size
            metadata['chunk_overlap_used'] = chunk_overlap
            metadata['chunking_mode'] = 'tokens' if self.use_tokens else 'characters'
            
            # Add chunk-level content hash for granular duplicate detection
            chunk_content_hash = Document._compute_content_hash(chunk_text)
            metadata['chunk_content_hash'] = chunk_content_hash
            
            chunk = DocumentChunk(
                content=chunk_text,
                metadata=metadata,
                chunk_id=None  # Will be assigned by vector store
            )
            chunks.append(chunk)
        
        return chunks
    
    def chunk_documents(self, documents: List[Document]) -> List[DocumentChunk]:
        """
        Split multiple documents into chunks.
        
        Args:
            documents: List of documents to chunk
            
        Returns:
            List of DocumentChunk objects
        """
        all_chunks = []
        for doc in documents:
            chunks = self.chunk_document(doc)
            all_chunks.extend(chunks)
        
        return all_chunks


if __name__ == "__main__":
    # Test the document processor
    processor = DocumentProcessor(chunk_size=200, chunk_overlap=20)
    
    # Test text splitting
    sample_text = """
    Artificial intelligence is the simulation of human intelligence by machines.
    Machine learning is a subset of AI that enables systems to learn from data.
    Deep learning uses neural networks with multiple layers.
    Natural language processing helps computers understand human language.
    """
    
    chunks = processor.split_text(sample_text)
    print(f"Created {len(chunks)} chunks from sample text")
    for i, chunk in enumerate(chunks):
        print(f"Chunk {i}: {chunk[:50]}...")
